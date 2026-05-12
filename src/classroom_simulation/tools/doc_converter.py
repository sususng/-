"""
文档转换工具模块

本模块提供文档格式转换工具，用于将Office文档和PDF转换为Markdown格式：
- DocConvertTool: 单个文档转换工具（支持.docx, .doc, .pdf, .pptx, .xlsx）
- BatchConvertTool: 批量文档转换工具

作者: AI备课系统
版本: 2.0
"""

from crewai.tools import BaseTool
from pydantic import Field
from typing import Type
from pathlib import Path
from datetime import datetime


class DocConvertInput(BaseTool):
    """
    文档转换输入参数类

    定义文档转换时需要的输入参数：
    - source_path: 源文件路径
    - output_path: 输出Markdown文件路径（可选，自动命名）
    """
    source_path: str = Field(..., description="源文件路径（.docx/.doc/.pdf/.pptx/.xlsx）")
    output_path: str = Field(default="", description="输出Markdown文件路径，默认自动命名")


class DocConvertTool(BaseTool):
    """
    文档转换工具

    功能：将Word、PDF、PPT、Excel文档转换为Markdown格式，便于AI处理

    支持的格式：
    - .docx: Word 2007+文档
    - .doc: Word 97-2003文档
    - .pdf: PDF文档
    - .pptx: PowerPoint文档
    - .xlsx/.xls: Excel表格

    使用场景：
    - 将已有的教案文档转换为Markdown以便AI分析
    - 将PDF格式的教学资料转换为可编辑文本
    - 将PPT课件转换为Markdown以便AI理解内容

    示例：
        tool = DocConvertTool()
        result = tool._run("教案.docx", "教案.md")
    """

    name: str = "doc_convert"
    description: str = "将Word、PDF、PPT、Excel文档转换为Markdown格式，便于AI处理"
    args_schema: Type[DocConvertInput] = DocConvertInput

    def _run(self, source_path: str, output_path: str = "") -> str:
        """
        转换文档

        参数:
            source_path: str - 源文件路径
            output_path: str - 输出Markdown文件路径（可选，自动命名）

        返回:
            str: 操作结果，如果成功则返回成功信息，失败则返回错误信息
        """
        try:
            source = Path(source_path)
            if not source.exists():
                return f"错误：源文件不存在 - {source_path}"

            ext = source.suffix.lower()

            if ext == ".docx":
                return self._convert_docx(source, output_path)
            elif ext == ".doc":
                return self._convert_doc(source, output_path)
            elif ext == ".pdf":
                return self._convert_pdf(source, output_path)
            elif ext == ".pptx":
                return self._convert_pptx(source, output_path)
            elif ext in [".xlsx", ".xls"]:
                return self._convert_xlsx(source, output_path)
            else:
                return f"错误：不支持的格式 - {ext}，支持：.docx, .doc, .pdf, .pptx, .xlsx"

        except ImportError as e:
            return f"错误：缺少依赖库 - {str(e)}，请安装 python-docx, pdfplumber, python-pptx, openpyxl"
        except Exception as e:
            return f"转换出错：{str(e)}"

    def _convert_docx(self, source: Path, output_path: str) -> str:
        """
        转换Word 2007+文档

        从.docx文件中提取文本和表格，转换为Markdown格式。

        参数:
            source: Path - 源文件路径对象
            output_path: str - 输出文件路径

        返回:
            str: 操作结果
        """
        from docx import Document

        doc = Document(source)
        lines = []

        lines.append("---")
        lines.append(f"source: {source}")
        lines.append(f"converted: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append("format: docx")
        lines.append("---")
        lines.append("")
        lines.append(f"# {source.stem}")
        lines.append("")

        for para in doc.paragraphs:
            if para.text.strip():
                lines.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                lines.append("| " + " | ".join(row_data) + " |")
            lines.append("")

        content = "\n".join(lines)

        if not output_path:
            output_path = str(source.with_suffix(".md"))

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(content, encoding="utf-8")

        return f"成功转换：{source.name} → {output.name}"

    def _convert_doc(self, source: Path, output_path: str) -> str:
        """
        转换Word 97-2003文档

        从.doc文件中提取文本，转换为Markdown格式。

        参数:
            source: Path - 源文件路径对象
            output_path: str - 输出文件路径

        返回:
            str: 操作结果
        """
        try:
            import olefile

            ole = olefile.OleFileIO(str(source))
            if "WordDocument" in ole.listdir():
                stream = ole.openstream("WordDocument").read()
                text = self._extract_word_text(stream)

                lines = []
                lines.append("---")
                lines.append(f"source: {source}")
                lines.append(f"converted: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                lines.append("format: doc")
                lines.append("---")
                lines.append("")
                lines.append(f"# {source.stem}")
                lines.append("")
                lines.append(text)

                content = "\n".join(lines)

                if not output_path:
                    output_path = str(source.with_suffix(".md"))

                output = Path(output_path)
                output.parent.mkdir(parents=True, exist_ok=True)
                output.write_text(content, encoding="utf-8")

                return f"成功转换：{source.name} → {output.name}"
            else:
                return f"警告：无法从 {source.name} 中提取Word文档内容"
        except Exception as e:
            return f"转换.doc文件失败：{str(e)}"

    def _extract_word_text(self, data: bytes) -> str:
        """
        提取Word文档文本

        从Word文档的二进制数据中提取可读文本。

        参数:
            data: bytes - Word文档二进制数据

        返回:
            str: 提取的文本内容
        """
        try:
            text = data.decode("utf-8", errors="ignore")
            text = "".join(c if ord(c) >= 32 or c in "\n\r\t" else " " for c in text)
            return text[:5000]
        except Exception:
            return "无法提取文本内容"

    def _convert_pdf(self, source: Path, output_path: str) -> str:
        """
        转换PDF文档

        从PDF文件中提取文本和表格，转换为Markdown格式。

        参数:
            source: Path - 源文件路径对象
            output_path: str - 输出文件路径

        返回:
            str: 操作结果
        """
        import pdfplumber

        all_text = []

        with pdfplumber.open(source) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    all_text.append(f"## 第{i}页\n\n{page_text}")

                tables = page.extract_tables()
                for table in tables:
                    table_text = "\n".join(["| " + " | ".join([str(cell) if cell else "" for cell in row]) + " |" for row in table])
                    all_text.append(f"\n### 第{i}页表格\n\n{table_text}\n")

        lines = []
        lines.append("---")
        lines.append(f"source: {source}")
        lines.append(f"converted: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append("format: pdf")
        lines.append("---")
        lines.append("")
        lines.append(f"# {source.stem}")
        lines.append("")
        lines.append("\n\n".join(all_text))

        content = "\n".join(lines)

        if not output_path:
            output_path = str(source.with_suffix(".md"))

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(content, encoding="utf-8")

        return f"成功转换：{source.name} → {output.name}"

    def _convert_pptx(self, source: Path, output_path: str) -> str:
        """
        转换PPT文档

        从PowerPoint文件中提取每页文本，转换为Markdown格式。

        参数:
            source: Path - 源文件路径对象
            output_path: str - 输出文件路径

        返回:
            str: 操作结果
        """
        from pptx import Presentation

        prs = Presentation(source)
        all_slides = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"## 第{i}页\n"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)

            if len(slide_text) > 1:
                all_slides.append("\n".join(slide_text))

        lines = []
        lines.append("---")
        lines.append(f"source: {source}")
        lines.append(f"converted: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append("format: pptx")
        lines.append("---")
        lines.append("")
        lines.append(f"# {source.stem}")
        lines.append("")
        lines.append("\n\n---\n\n".join(all_slides))

        content = "\n".join(lines)

        if not output_path:
            output_path = str(source.with_suffix(".md"))

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(content, encoding="utf-8")

        return f"成功转换：{source.name} → {output.name}"

    def _convert_xlsx(self, source: Path, output_path: str) -> str:
        """
        转换Excel文档

        从Excel文件中提取所有工作表数据，转换为Markdown表格格式。

        参数:
            source: Path - 源文件路径对象
            output_path: str - 输出文件路径

        返回:
            str: 操作结果
        """
        from openpyxl import load_workbook

        wb = load_workbook(source, data_only=True)
        all_sheets = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = [f"## 工作表: {sheet_name}\n"]

            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(row_data):
                    sheet_text.append("| " + " | ".join(row_data) + " |")

            all_sheets.append("\n".join(sheet_text))

        lines = []
        lines.append("---")
        lines.append(f"source: {source}")
        lines.append(f"converted: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append("format: xlsx")
        lines.append("---")
        lines.append("")
        lines.append(f"# {source.stem}")
        lines.append("")
        lines.append("\n\n---\n\n".join(all_sheets))

        content = "\n".join(lines)

        if not output_path:
            output_path = str(source.with_suffix(".md"))

        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(content, encoding="utf-8")

        return f"成功转换：{source.name} → {output.name}"


class BatchConvertInput(BaseTool):
    """
    批量转换输入参数类

    定义批量转换时需要的输入参数：
    - source_dir: 源文件目录路径
    - output_dir: 输出目录路径（可选，默认与源目录相同）
    """
    source_dir: str = Field(..., description="源文件目录路径")
    output_dir: str = Field(default="", description="输出目录路径，默认与源目录相同")


class BatchConvertTool(BaseTool):
    """
    批量文档转换工具

    功能：批量将目录下所有支持的文档（Word、PDF、PPT、Excel）转换为Markdown

    使用场景：
    - 批量转换教学资料文件夹中的所有文档
    - 统一处理多个教案文件
    - 批量提取PDF教材内容

    示例：
        tool = BatchConvertTool()
        result = tool._run("教学资料", "输出")
    """

    name: str = "batch_convert"
    description: str = "批量将目录下所有Word、PDF、PPT、Excel文档转换为Markdown"
    args_schema: Type[BatchConvertInput] = BatchConvertInput

    def _run(self, source_dir: str, output_dir: str = "") -> str:
        """
        批量转换

        参数:
            source_dir: str - 源文件目录路径
            output_dir: str - 输出目录路径（可选，默认与源目录相同）

        返回:
            str: 操作结果汇总，包含成功数量、失败数量和详细结果
        """
        try:
            source_path = Path(source_dir)
            if not source_path.exists():
                return f"错误：源目录不存在 - {source_dir}"

            if not output_dir:
                output_dir = source_dir
            output_path = Path(output_dir)

            supported_exts = [".docx", ".doc", ".pdf", ".pptx", ".xlsx", ".xls"]
            results = []
            success_count = 0
            fail_count = 0

            converter = DocConvertTool()

            for ext in supported_exts:
                for file in source_path.rglob(f"*{ext}"):
                    try:
                        result = converter._run(str(file), str(output_path / file.with_suffix(".md").name))
                        if "成功" in result:
                            success_count += 1
                            results.append(f"✅ {result}")
                        else:
                            fail_count += 1
                            results.append(f"❌ {file.name}: {result}")
                    except Exception as e:
                        fail_count += 1
                        results.append(f"❌ {file.name}: {str(e)}")

            summary = f"\n\n转换完成：成功 {success_count} 个，失败 {fail_count} 个"
            return summary + "\n\n" + "\n".join(results)

        except Exception as e:
            return f"批量转换出错：{str(e)}"